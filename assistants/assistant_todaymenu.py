
import json
from openai import OpenAI
import os

from io import BufferedReader
from typing import Optional
#from fastapi import UploadFile
import mimetypes
from PyPDF2 import PdfReader
import docx2txt
import csv
import pptx

client = OpenAI()

todaymenu_path="../files/menu.csv"
class TodayMenuAssistant:
    def __init__(self, client):
        self.client = client
        self.message_file_id = None
        self.file_path = todaymenu_path
        self.upload_file()

    def upload_file(self):
        current_dir = os.path.dirname(os.path.abspath(__file__))
        self.file_path = os.path.join(current_dir, todaymenu_path)

        message_file = self.client.files.create(
            file=open(self.file_path, "rb"), purpose="assistants"
        )
        print("Uploaded file: " + message_file.id)
        self.message_file_id = message_file.id

    def initModule(client, assistant):
        vector_store = client.beta.vector_stores.create(name="Today Menu")
        current_dir = os.path.dirname(os.path.abspath(__file__))
        file_path = os.path.join(current_dir, "..", "files", "menu.csv")
        file_paths = [file_path]
        file_streams = [open(path, "rb") for path in file_paths]

        file_batch = client.beta.vector_stores.file_batches.upload_and_poll(
            vector_store_id=vector_store.id, files=file_streams
        )
        print(file_batch.status)
        print(file_batch.file_counts)

        assistant = client.beta.assistants.update(
            assistant_id=assistant.id,
            tool_resources={"file_search": {"vector_store_ids": [vector_store.id]}},
        )

        return assistant

    def createFunctions_todaymenu(self, client):
        return [
            {
                "type": "function",
                "function": {
                    "name": "get_today_menu",
                    "description": "Get the today menu",
                    "parameters": {
                        "type": "object",
                        "properties": {
                            "preference": {
                                "type": "string",
                                "description": "If they mention any preference"
                            }                        
                        },
                        "required": [
                        ]
                    }
                }
            }#,
         #{"type": "file_search"}
        ]


    def handleToolCall(self, function_name, arguments=None, client=None):

        if function_name == "get_today_menu":
            ##print("Uploading file...")

            #cont = self.extract_text_from_filepath(self.file_path, mimetype="csv")
            cont = ""
            with open(self.file_path, "rb") as file:
                cont = file.read().decode("utf-8")

            # Example response
            menu_data = {
                "content": cont,#"Find the Today's menu in the attached file",
                "attachments": [{"file_id": self.message_file_id, "tools": [{"type": "file_search"}]}]
            }
            return json.dumps(menu_data)

        return None

    def extract_text_from_filepath(self, filepath: str, mimetype: Optional[str] = None) -> str:
        """Return the text content of a file given its filepath."""
        
        print(filepath)

        if mimetype is None:
            # Get the mimetype of the file based on its extension
            mimetype, _ = mimetypes.guess_type(filepath)

        if not mimetype:
            if filepath.endswith(".md"):
                mimetype = "text/markdown"
            else:
                raise Exception("Unsupported file type")

        try:
            with open(filepath, "rb") as file:
                extracted_text = self.extract_text_from_file(file, mimetype)
        except Exception as e:
            #logger.error(e)
            raise e

        return extracted_text


    def extract_text_from_file(self, file: BufferedReader, mimetype: str) -> str:
        if mimetype == "application/pdf":
            # Extract text from pdf using PyPDF2
            reader = PdfReader(file)
            extracted_text = " ".join([page.extract_text() for page in reader.pages])
        elif mimetype == "text/plain" or mimetype == "text/markdown":
            # Read text from plain text file
            extracted_text = file.read().decode("utf-8")
        elif (
            mimetype
            == "application/vnd.openxmlformats-officedocument.wordprocessingml.document"
        ):
            # Extract text from docx using docx2txt
            extracted_text = docx2txt.process(file)
        elif mimetype == "text/csv":
            # Extract text from csv using csv module
            extracted_text = ""
            decoded_buffer = (line.decode("utf-8") for line in file)
            reader = csv.reader(decoded_buffer)
            for row in reader:
                extracted_text += " ".join(row) + "\n"
        elif (
            mimetype
            == "application/vnd.openxmlformats-officedocument.presentationml.presentation"
        ):
            # Extract text from pptx using python-pptx
            extracted_text = ""
            presentation = pptx.Presentation(file)
            for slide in presentation.slides:
                for shape in slide.shapes:
                    if shape.has_text_frame:
                        for paragraph in shape.text_frame.paragraphs:
                            for run in paragraph.runs:
                                extracted_text += run.text + " "
                        extracted_text += "\n"
        else:
            # Unsupported file type
            raise ValueError("Unsupported file type: {}".format(mimetype))

        return extracted_text


    # # Extract text from a file based on its mimetype
    # async def extract_text_from_form_file(self, file: UploadFile):
    #     """Return the text content of a file."""
    #     # get the file body from the upload file object
    #     mimetype = file.content_type
    #     # logger.info(f"mimetype: {mimetype}")
    #     # logger.info(f"file.file: {file.file}")
    #     # logger.info("file: ", file)

    #     file_stream = await file.read()

    #     temp_file_path = "/tmp/temp_file"

    #     # write the file to a temporary location
    #     with open(temp_file_path, "wb") as f:
    #         f.write(file_stream)

    #     try:
    #         extracted_text =self.extract_text_from_filepath(temp_file_path, mimetype)
    #     except Exception as e:
    #         # logger.error(e)
    #         os.remove(temp_file_path)
    #         raise e

    #     # remove file from temp location
    #     os.remove(temp_file_path)

    #     return extracted_text

today_menu_assistant = TodayMenuAssistant(client)

    #if function_name == "get_today_menu":
    #     menu_data = {
    #         "menu": "Today's menu: Pizza, Pasta, Salad"
    #     }
    # message_file = client.files.create(
    #     file=open("files/menu.pdf", "rb"), purpose="assistants"
    # )
    # print("attaching" + message_file_id)

    # return  json.dumps({
    #     #"content": json.dumps(function_response),
    #     "attachments": [
    #         { "file_id": message_file_id, "tools": [{"type": "file_search"}] }
    #     ]})
    
    #return json.dumps(menu_data)

    #return False


    # if(function_name in available_weather_functions):
    #     function_to_call =  available_weather_functions[function_name]
    #     function_args = json.loads(arguments)
    #     print(function_args)
    #     # function_response = function_to_call(
    #     #    #location=function_args.get("location"),
    #     #    # unit=function_args.get("unit", "c"),
    #     # )

    #     return  json.dumps({
    #         #"content": json.dumps(function_response),
    #         "attachments": [{
    #             "message_fileid": message_file_id
    #         }]})
        
    
    
    # return False

# def get_menu():
#     return {}

# available_weather_functions = {
#         "get_today_menu": get_menu,
#     }




   #print('weather tool')
      # only one function in this example, but you can have multiple

    # vector_store = client.beta.vector_stores.create(name="Today Menu")
    # file_paths = ["../files/menu.pdf"]
    # file_streams = [open(path, "rb") for path in file_paths]

    # file_batch = client.beta.vector_stores.file_batches.upload_and_poll(
    #     vector_store_id=vector_store.id, files=file_streams
    # )
    # print(file_batch.status)
    # print(file_batch.file_counts)

    # assistant = client.beta.assistants.update(
    #     assistant_id=assistant.id,
    #     tool_resources={"file_search": {"vector_store_ids": [vector_store.id]}},
    # )